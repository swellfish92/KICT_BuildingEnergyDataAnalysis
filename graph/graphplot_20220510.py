import scipy as sp
import scipy.stats
import pymysql
import pandas as pd
from functions import *
import numpy as np
from sklearn.linear_model import LinearRegression
import matplotlib.transforms as trf
from pptx import Presentation
from pptx.util import Inches



month_arr = ['01', '02', '03', '04', '05', '06', '07', '08', '09', '10', '11', '12']
year_arr = ['2014', '2015', '2016', '2017', '2018', '2019', '2020']
year_arr_int = [2014, 2015, 2016, 2017, 2018, 2019, 2020]



def print_annual(basedir, set_pk, pk_arr):
    original_data = pd.read_csv(basedir)
    print(original_data)
    original_data.set_index('mgmBldrgstPk', drop=True, inplace=True)
    for idx, row in original_data.iterrows():
        if idx not in pk_arr:
            original_data = original_data.drop(idx)
    print(original_data)

    for idx, row in original_data.iterrows():
        print(row)


        temp_arr = []
        for year in year_arr:
            carbon_val = 0
            for month in month_arr:
                carbon_val = carbon_val + row['Carbon_'+year+month]
            print(carbon_val)
            temp_arr.append(carbon_val/(row['TOTAREA']))  # 단위면적당 tCO2
            # temp_arr.append(carbon_val/1000)    # 총 tCO2
        #
        # temp_arr = []
        # for year in year_arr:
        #     carbon_val = 0
        #     for month in month_arr:
        #         carbon_val = carbon_val + row['Carbon_' + year + month]
        #     if year_arr.index(year) == 0:
        #         # temp_arr.append(carbon_val/(row['TOTAREA']))  # 단위면적당 tCO2
        #         temp_arr.append(carbon_val / 1000)  # 총 tCO2
        #     else:
        #         # temp_arr.append(carbon_val/(row['TOTAREA']))  # 단위면적당 tCO2
        #         temp_arr.append(carbon_val / 1000)  # 총 tCO2
        # print(temp_arr)
        # temp_arr_2 = []
        # # 도중에 생긴 건물이면 2014년 기준이 안되니까 제거함
        # if temp_arr[0] == 0:
        #     continue
        # # 기타 극단값 필터들
        # if temp_arr[6]/temp_arr[0] > 3:
        #     continue
        # if temp_arr[6]/temp_arr[0] < 0.8:
        #     continue
        #
        # for item in temp_arr:
        #     temp_arr_2.append(item/temp_arr[0])
        # temp_arr = temp_arr_2
        # # temp_arr.append(carbon_val/1000)    # 총 tCO2

        print(row['Univ_name'])
        if idx == set_pk:
            plt.plot(year_arr_int, temp_arr, color='red', marker='o', linewidth=2, markersize=12, label=row['Univ_name'], alpha=0.8)
            for i in range(len(year_arr_int)):
                plt.text(year_arr_int[i], temp_arr[i]+1, str(round(temp_arr[i], 2)))
        else:
            plt.plot(year_arr_int, temp_arr, color='lightgray', marker='o', linewidth=1, markersize=6, alpha=0.8)
            plt.text(year_arr_int[6]+0.05, temp_arr[6], row['Univ_name'])
    plt.legend()
    plt.show()



# target_pk = '11680-410' # 강남세브란스병원
# target_pk_2 = '11740-435' # 중앙보훈병원
# target_pk = '11290-1275' # 고려대학교
target_pk = '11710-228' # 서울아산병원
univ_arr = ['11110-100183537', '11110-12058', '11110-462', '11110-647', '11140-14986', '11140-176', '11170-171', '11200-1225', '11215-18377', '11215-23070', '11230-29344', '11230-34287', '11230-388', '11230-632', '11260-11032', '11260-25373', '11290-518', '11305-19037', '11320-280', '11350-14660', '11350-503', '11350-9982', '11380-100290794', '11380-14367', '11410-100185809', '11410-21340', '11470-100186127', '11470-100191796', '11470-7821', '11500-100214588', '11500-100318672', '11500-6410', '11530-5576', '11530-7681', '11545-12486', '11560-100194494', '11560-100251089', '11560-1390', '11560-33312', '11560-507', '11560-59', '11560-6251', '11590-100194692', '11590-14698', '11620-20883', '11650-323', '11680-410', '11680-74', '11680-88', '11710-228', '11710-485', '11740-435', '11740-4945', '11740-544']
basedir = 'C:/Users/user/Downloads/병원데이터시트_20220517_v2.csv'
print_annual(basedir, target_pk, univ_arr)
# print_annual(basedir, target_pk_2, univ_arr)



raise IOError




def get_median_v2(data):
    data = sorted(data)

    centerIndex = len(data) // 2
    return (data[centerIndex] + data[-centerIndex - 1]) / 2

# data = pd.read_csv('수작업시발.csv')
# print(data)
# print(data['난방기울기'].median())
# print(data['냉방기울기'].median())
# print(data['난방기울기'].median())
#
# raise IOError

original_data = pd.read_csv('C:/Users/user/Downloads/병원데이터시트_20220517_v2.csv')
basedata = original_data[['mgmBldrgstPk', 'TOTAREA', 'Carbon_201401', 'Carbon_201402', 'Carbon_201403', 'Carbon_201404', 'Carbon_201405', 'Carbon_201406', 'Carbon_201407', 'Carbon_201408', 'Carbon_201409', 'Carbon_201410', 'Carbon_201411', 'Carbon_201412', 'Carbon_201501', 'Carbon_201502', 'Carbon_201503', 'Carbon_201504', 'Carbon_201505', 'Carbon_201506', 'Carbon_201507', 'Carbon_201508', 'Carbon_201509', 'Carbon_201510', 'Carbon_201511', 'Carbon_201512', 'Carbon_201601', 'Carbon_201602', 'Carbon_201603', 'Carbon_201604', 'Carbon_201605', 'Carbon_201606', 'Carbon_201607', 'Carbon_201608', 'Carbon_201609', 'Carbon_201610', 'Carbon_201611', 'Carbon_201612', 'Carbon_201701', 'Carbon_201702', 'Carbon_201703', 'Carbon_201704', 'Carbon_201705', 'Carbon_201706', 'Carbon_201707', 'Carbon_201708', 'Carbon_201709', 'Carbon_201710', 'Carbon_201711', 'Carbon_201712', 'Carbon_201801', 'Carbon_201802', 'Carbon_201803', 'Carbon_201804', 'Carbon_201805', 'Carbon_201806', 'Carbon_201807', 'Carbon_201808', 'Carbon_201809', 'Carbon_201810', 'Carbon_201811', 'Carbon_201812', 'Carbon_201901', 'Carbon_201902', 'Carbon_201903', 'Carbon_201904', 'Carbon_201905', 'Carbon_201906', 'Carbon_201907', 'Carbon_201908', 'Carbon_201909', 'Carbon_201910', 'Carbon_201911', 'Carbon_201912', 'Carbon_202001', 'Carbon_202002', 'Carbon_202003', 'Carbon_202004', 'Carbon_202005', 'Carbon_202006', 'Carbon_202007', 'Carbon_202008', 'Carbon_202009', 'Carbon_202010', 'Carbon_202011', 'Carbon_202012']]    # , 'Carbon_202101', 'Carbon_202102', 'Carbon_202103', 'Carbon_202104', 'Carbon_202107', 'Carbon_202108', 'Carbon_202109', 'Carbon_202110', 'Carbon_202111', 'Carbon_202112']]
univname_data = original_data[['mgmBldrgstPk', 'Univ_name']]
univname_data.set_index('mgmBldrgstPk', drop=True, inplace=True)
basedata.set_index('mgmBldrgstPk', drop=True, inplace=True)
temperature_data = pd.read_csv('C:/Users/user/Downloads/ta_20220510173204.csv')
print(basedata)

data_arr = ['Carbon_201401', 'Carbon_201402', 'Carbon_201403', 'Carbon_201404', 'Carbon_201405', 'Carbon_201406', 'Carbon_201407', 'Carbon_201408', 'Carbon_201409', 'Carbon_201410', 'Carbon_201411', 'Carbon_201412', 'Carbon_201501', 'Carbon_201502', 'Carbon_201503', 'Carbon_201504', 'Carbon_201505', 'Carbon_201506', 'Carbon_201507', 'Carbon_201508', 'Carbon_201509', 'Carbon_201510', 'Carbon_201511', 'Carbon_201512', 'Carbon_201601', 'Carbon_201602', 'Carbon_201603', 'Carbon_201604', 'Carbon_201605', 'Carbon_201606', 'Carbon_201607', 'Carbon_201608', 'Carbon_201609', 'Carbon_201610', 'Carbon_201611', 'Carbon_201612', 'Carbon_201701', 'Carbon_201702', 'Carbon_201703', 'Carbon_201704', 'Carbon_201705', 'Carbon_201706', 'Carbon_201707', 'Carbon_201708', 'Carbon_201709', 'Carbon_201710', 'Carbon_201711', 'Carbon_201712', 'Carbon_201801', 'Carbon_201802', 'Carbon_201803', 'Carbon_201804', 'Carbon_201805', 'Carbon_201806', 'Carbon_201807', 'Carbon_201808', 'Carbon_201809', 'Carbon_201810', 'Carbon_201811', 'Carbon_201812', 'Carbon_201901', 'Carbon_201902', 'Carbon_201903', 'Carbon_201904', 'Carbon_201905', 'Carbon_201906', 'Carbon_201907', 'Carbon_201908', 'Carbon_201909', 'Carbon_201910', 'Carbon_201911', 'Carbon_201912', 'Carbon_202001', 'Carbon_202002', 'Carbon_202003', 'Carbon_202004', 'Carbon_202005', 'Carbon_202006', 'Carbon_202007', 'Carbon_202008', 'Carbon_202009', 'Carbon_202010', 'Carbon_202011', 'Carbon_202012'] # , 'Carbon_202101', 'Carbon_202102', 'Carbon_202103', 'Carbon_202104',  'Carbon_202107', 'Carbon_202108', 'Carbon_202109', 'Carbon_202110', 'Carbon_202111', 'Carbon_202112']

# 면적별로 교환해주는 코드임.
for idx in data_arr:
    basedata[idx] = basedata[idx] / basedata['TOTAREA']

basedata = basedata.drop(labels=['TOTAREA'], axis=1)
print(basedata)

target_pk_arr = ['11680-410'] # 강남세브란스병원
#target_pk_arr = ['11440-310']
target_pk_arr = ['11740-435'] # 중앙보훈병원
#target_pk_arr = ['11410-100185809'] # 연대

prs = Presentation()

# 병원비교용 별도코드

# hospital_arr = ['서울대학교병원', '삼성서울병원', '재단법인아산사회복지재단 서울아산병원']
# color_dict = {'서울대학교병원':'blue', '삼성서울병원':'red', '재단법인아산사회복지재단 서울아산병원':'green'}
# temp_dict = {}
# for idx in basedata.index.tolist():
#
#     carbon_arr = basedata.loc[idx].values.tolist()
#     # temp_arr = [-0.7, 1.9, 7.9, 14, 18.9, 23.1, 26.1, 25.2, 22.1, 15.6, 9, -2.9, -0.9, 1, 6.3, 13.3, 18.9, 23.6, 25.8, 26.3, 22.4, 15.5, 8.9, 1.6, -3.2, 0.2, 7, 14.1, 19.6, 23.6, 26.2, 28, 23.1, 16.1, 6.8, 1.2, -1.8, -0.2, 6.3, 13.9, 19.5, 23.3, 26.9, 25.9, 22.1, 16.4, 5.6, -1.9, -4, -1.6, 8.1, 13, 18.2, 23.1, 27.8, 28.8, 21.5, 13.1, 7.8, -0.6, -0.9, 1, 7.1, 12.1, 19.4, 22.5, 25.9, 27.2, 22.6, 16.4, 7.6, 1.4, 1.6, 2.5, 7.7, 11.1, 18, 23.9, 24.1, 26.5, 21.4, 14.3, 8, -0.3, -2.4, 2.7, 9, 14.2, 28.1, 25.9, 22.6, 15.6, 8.2, 0.6]
#     temp_arr = [-0.7, 1.9, 7.9, 14, 18.9, 23.1, 26.1, 25.2, 22.1, 15.6, 9, -2.9, -0.9, 1, 6.3, 13.3, 18.9, 23.6, 25.8, 26.3, 22.4, 15.5, 8.9, 1.6, -3.2, 0.2, 7, 14.1, 19.6, 23.6, 26.2, 28, 23.1, 16.1, 6.8, 1.2, -1.8, -0.2, 6.3, 13.9, 19.5, 23.3, 26.9, 25.9, 22.1, 16.4, 5.6, -1.9, -4, -1.6, 8.1, 13, 18.2, 23.1, 27.8, 28.8, 21.5, 13.1, 7.8, -0.6, -0.9, 1, 7.1, 12.1, 19.4, 22.5, 25.9, 27.2, 22.6, 16.4, 7.6, 1.4, 1.6, 2.5, 7.7, 11.1, 18, 23.9, 24.1, 26.5, 21.4, 14.3, 8, -0.3]
#
#     univ_name = univname_data.loc[idx]['Univ_name']
#     print(univ_name)
#
#     if univ_name in hospital_arr:
#         plt.scatter(temp_arr, carbon_arr, c=color_dict[univ_name], s=10, label=univ_name)
#
#         slope_arr = []
#
#         # 선 그리기 (난방선)
#         heating_carbon_arr = []
#         heating_temp_arr = []
#         for i in range(len(temp_arr)):
#             if temp_arr[i] < 10:
#                 heating_carbon_arr.append(carbon_arr[i])
#                 heating_temp_arr.append(temp_arr[i])
#         line_reg_heating = LinearRegression()
#         line_reg_heating.fit(np.reshape(heating_temp_arr, (-1, 1)), np.reshape(heating_carbon_arr, (-1, 1)))
#
#         #
#
#         slope_arr.append(line_reg_heating.coef_[0][0])
#
#         # 선 그리기 (냉방선)
#         cooling_carbon_arr = []
#         cooling_temp_arr = []
#         for i in range(len(temp_arr)):
#             if temp_arr[i] >= 20:
#                 cooling_carbon_arr.append(carbon_arr[i])
#                 cooling_temp_arr.append(temp_arr[i])
#         line_reg_cooling = LinearRegression()
#         line_reg_cooling.fit(np.reshape(cooling_temp_arr, (-1, 1)), np.reshape(cooling_carbon_arr, (-1, 1)))
#         slope_arr.append(line_reg_cooling.coef_[0][0])
#
#         # 기저사용량 뽑기
#         base_usage_arr = []
#         for i in range(len(temp_arr)):
#             if temp_arr[i] < 17.5 and temp_arr[i] >= 12.5:
#                 base_usage_arr.append(carbon_arr[i])
#         slope_arr.append(get_median_v2(base_usage_arr))
#
#         plt.plot(heating_temp_arr, line_reg_heating.predict(np.reshape(heating_temp_arr, (-1, 1))), color=color_dict[univ_name])
#         plt.plot(cooling_temp_arr, line_reg_cooling.predict(np.reshape(cooling_temp_arr, (-1, 1))), color=color_dict[univ_name])
#         plt.plot([5, 25], [get_median_v2(base_usage_arr), get_median_v2(base_usage_arr)], color=color_dict[univ_name], linewidth=1,
#                  linestyle='--')
#         print(slope_arr)
#
#         plt.text(-8, line_reg_heating.predict([[-4]]), '난방 : ' + str(round(slope_arr[0], 4)), color=color_dict[univ_name])
#         plt.text(29.5, line_reg_cooling.predict([[27.5]]), '냉방 : ' + str(round(slope_arr[1], 4)), color=color_dict[univ_name])
#         plt.text(27.5, get_median_v2(base_usage_arr), '기저 : ' + str(round(slope_arr[2], 4)), color=color_dict[univ_name])
#         temp_dict[idx] = slope_arr
#         plt.title(univ_name + '//' + idx)
#         plt.xlim(-10, 40)
#
#
# plt.legend()
# plt.show()
# raise IOError

temp_dict = {}
for idx in basedata.index.tolist():

    carbon_arr = basedata.loc[idx].values.tolist()
    # temp_arr = [-0.7, 1.9, 7.9, 14, 18.9, 23.1, 26.1, 25.2, 22.1, 15.6, 9, -2.9, -0.9, 1, 6.3, 13.3, 18.9, 23.6, 25.8, 26.3, 22.4, 15.5, 8.9, 1.6, -3.2, 0.2, 7, 14.1, 19.6, 23.6, 26.2, 28, 23.1, 16.1, 6.8, 1.2, -1.8, -0.2, 6.3, 13.9, 19.5, 23.3, 26.9, 25.9, 22.1, 16.4, 5.6, -1.9, -4, -1.6, 8.1, 13, 18.2, 23.1, 27.8, 28.8, 21.5, 13.1, 7.8, -0.6, -0.9, 1, 7.1, 12.1, 19.4, 22.5, 25.9, 27.2, 22.6, 16.4, 7.6, 1.4, 1.6, 2.5, 7.7, 11.1, 18, 23.9, 24.1, 26.5, 21.4, 14.3, 8, -0.3, -2.4, 2.7, 9, 14.2, 28.1, 25.9, 22.6, 15.6, 8.2, 0.6]
    temp_arr = [-0.7, 1.9, 7.9, 14, 18.9, 23.1, 26.1, 25.2, 22.1, 15.6, 9, -2.9, -0.9, 1, 6.3, 13.3, 18.9, 23.6, 25.8, 26.3, 22.4, 15.5, 8.9, 1.6, -3.2, 0.2, 7, 14.1, 19.6, 23.6, 26.2, 28, 23.1, 16.1, 6.8, 1.2, -1.8, -0.2, 6.3, 13.9, 19.5, 23.3, 26.9, 25.9, 22.1, 16.4, 5.6, -1.9, -4, -1.6, 8.1, 13, 18.2, 23.1, 27.8, 28.8, 21.5, 13.1, 7.8, -0.6, -0.9, 1, 7.1, 12.1, 19.4, 22.5, 25.9, 27.2, 22.6, 16.4, 7.6, 1.4, 1.6, 2.5, 7.7, 11.1, 18, 23.9, 24.1, 26.5, 21.4, 14.3, 8, -0.3]

    univ_name = univname_data.loc[idx]['Univ_name']

    plt.scatter(temp_arr, carbon_arr, c='blue', s=1.5)
    # # 점 뿌리기
    # if idx in target_pk_arr:
    #     plt.scatter(temp_arr, carbon_arr, c='red', label=idx)
    # else:
    #     plt.scatter(temp_arr, carbon_arr, c='blue', s=1.5)

    slope_arr = []

    # 선 그리기 (난방선)
    heating_carbon_arr = []
    heating_temp_arr = []
    for i in range(len(temp_arr)):
        if temp_arr[i] < 10:
            heating_carbon_arr.append(carbon_arr[i])
            heating_temp_arr.append(temp_arr[i])
    line_reg_heating = LinearRegression()
    line_reg_heating.fit(np.reshape(heating_temp_arr, (-1,1)), np.reshape(heating_carbon_arr, (-1,1)))

    #


    slope_arr.append(line_reg_heating.coef_[0][0])
    
    # 선 그리기 (냉방선)
    cooling_carbon_arr = []
    cooling_temp_arr = []
    for i in range(len(temp_arr)):
        if temp_arr[i] >= 20:
            cooling_carbon_arr.append(carbon_arr[i])
            cooling_temp_arr.append(temp_arr[i])
    line_reg_cooling = LinearRegression()
    line_reg_cooling.fit(np.reshape(cooling_temp_arr, (-1,1)), np.reshape(cooling_carbon_arr, (-1,1)))
    slope_arr.append(line_reg_cooling.coef_[0][0])

    # 기저사용량 뽑기
    base_usage_arr = []
    for i in range(len(temp_arr)):
        if temp_arr[i] < 17.5 and temp_arr[i] >= 12.5:
            base_usage_arr.append(carbon_arr[i])
    slope_arr.append(get_median_v2(base_usage_arr))

    plt.plot(heating_temp_arr, line_reg_heating.predict(np.reshape(heating_temp_arr, (-1, 1))), color='black')
    plt.plot(cooling_temp_arr, line_reg_cooling.predict(np.reshape(cooling_temp_arr, (-1, 1))), color='black')
    plt.plot([5, 25], [get_median_v2(base_usage_arr), get_median_v2(base_usage_arr)], color='red', linewidth=1, linestyle = '--')
    print(slope_arr)

    plt.text(1, line_reg_heating.predict([[0]]), '난방경사 : ' + str(round(slope_arr[0], 4)))
    plt.text(26.5, line_reg_cooling.predict([[25]]), '냉방경사 : ' + str(round(slope_arr[1], 4)))
    plt.text(25.5, get_median_v2(base_usage_arr), '기저량 : ' + str(round(slope_arr[2], 4)))
    temp_dict[idx] = slope_arr
    plt.title(univ_name + '//' + idx)
    plt.xlim(-10, 40)
    plt.savefig('result_pics/' + univ_name + '(' + idx + ').png', dpi=200)

    # 저장그림 PPT에 넣기
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.8)
    top = Inches(1)
    width = Inches(8)
    height = Inches(6)
    pic = slide.shapes.add_picture('result_pics/' + univ_name + '(' + idx + ').png', left, top, width=width, height=height)

    # plt.show()
    # plt.savefig('result_pics/' + univ_name + '(' + idx + ').png', dpi=200)
    plt.clf()

prs.save('병원_result_divarea_raw.pptx')

# raise IOError

# 키(PK코드)별로 재정렬
res_dict = {
    'mgmBldrgstPk':[],
    '냉방기울기':[],
    '난방기울기':[],
    '기저사용량':[]
}
for key in temp_dict.keys():
    res_dict['mgmBldrgstPk'].append(key)
    res_dict['난방기울기'].append(temp_dict[key][0])
    res_dict['냉방기울기'].append(temp_dict[key][1])
    res_dict['기저사용량'].append(temp_dict[key][2])

print(res_dict)

append_data = pd.DataFrame(res_dict)
append_data.set_index('mgmBldrgstPk', drop=True, inplace=True)

append_data.to_csv('수작업_병원목록.csv', encoding='utf-8 sig')

original_data = original_data.join(append_data, how='left')
print(original_data)
original_data.to_csv('기울기붙인최종버전_병원_20220516.csv', encoding='utf-8 sig')
# plt.legend()
# plt.show()



# plt.scatter(data)
# plt.show()

